"""
Slide Parser - Extracts content from PDF/PPT slide decks.
Supports both text extraction and optional vision-based parsing for image-heavy slides.
"""

from pathlib import Path
from typing import List, Dict, Optional
import base64

import fitz
import pymupdf  # PyMuPDF for PDF parsing
from pptx import Presentation  # python-pptx for PowerPoint
from dataclasses import dataclass


@dataclass
class Slide:
    """Represents a single slide with its content."""
    slide_number: int
    title: str
    content: str
    notes: str = ""
    image_data: Optional[bytes] = None
    
    def to_dict(self) -> Dict:
        """Convert slide to dictionary."""
        return {
            "slide_number": self.slide_number,
            "title": self.title,
            "content": self.content,
            "notes": self.notes,
            "image_data": base64.b64encode(self.image_data).decode('utf-8') if self.image_data else None,
            "has_image": self.image_data is not None
        }
    
    @classmethod
    def from_dict(cls, data: Dict) -> 'Slide':
        """Create Slide from dictionary."""
        image_data = base64.b64decode(data['image_data']) if data.get('image_data') else None
        return cls(
            slide_number=data['slide_number'],
            title=data['title'],
            content=data['content'],
            notes=data.get('notes', ""),
            image_data=image_data
        )


class SlideParser:
    """Parse slide decks and extract structured content."""
    
    def __init__(self, use_vision: bool = False, zoom: float = 3.0):
        """
        Initialize slide parser.
        
        Args:
            use_vision: Whether to use vision models for image-heavy slides
            zoom: Zoom factor for rendering PDF pages as images (higher value = better quality)
        """
        self.use_vision = use_vision
        self.zoom = zoom

    def parse(self, file_path: Path, support_files: Optional[List[Path]] = None) -> List[Slide]:
        """
        Parse a slide deck file.
        
        Args:
            file_path: Path to the slide deck (PDF or PPT)
            support_files: Optional list of supporting documents
            
        Returns:
            List of Slide objects
        """
        file_extension = file_path.suffix.lower()
        
        if file_extension == '.pdf':
            return self._parse_pdf(file_path)
        elif file_extension in ['.pptx', '.ppt']:
            return self._parse_pptx(file_path)
        else:
            raise ValueError(f"Unsupported file format: {file_extension}")
    
    def _parse_pdf(self, file_path: Path) -> List[Slide]:
        """Parse PDF slide deck."""
        slides = []
        
        with pymupdf.open(file_path) as doc:
            for page_num, page in enumerate(doc, start=1):
                # Extract text
                text = page.get_text()
                
                # Try to extract title (first line or largest text)
                lines = text.split('\n')
                title = lines[0] if lines else f"Slide {page_num}"
                content = '\n'.join(lines[1:]) if len(lines) > 1 else text
                
                # Optionally extract page as image for vision models
                image_data = None
                if self.use_vision:
                    mat = fitz.Matrix(self.zoom, self.zoom)  # Use configurable zoom
                    pix = page.get_pixmap(matrix=mat)
                    image_data = pix.tobytes("png")
                
                slides.append(Slide(
                    slide_number=page_num,
                    title=title.strip(),
                    content=content.strip(),
                    image_data=image_data
                ))
        
        return slides
    
    def _parse_pptx(self, file_path: Path) -> List[Slide]:
        """Parse PowerPoint slide deck."""
        slides = []
        prs = Presentation(file_path)

        for slide_num, slide in enumerate(prs.slides, start=1):
            title = ""
            content_parts = []
            notes = ""
            
            # Extract title
            if slide.shapes.title:
                title = slide.shapes.title.text
            
            # Extract content from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    if shape != slide.shapes.title:
                        content_parts.append(shape.text)
            
            # Extract notes
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                if notes_slide.notes_text_frame:
                    notes = notes_slide.notes_text_frame.text
            
            content = '\n'.join(content_parts)
            
            # Optionally render slide as image for vision models
            image_data = None
            if self.use_vision:
                # Note: Rendering PPTX to image requires additional libraries
                # This is a placeholder for future implementation
                pass
            
            slides.append(Slide(
                slide_number=slide_num,
                title=title.strip() if title else f"Slide {slide_num}",
                content=content.strip(),
                notes=notes.strip(),
                image_data=image_data
            ))
        
        return slides
    
    def extract_with_vision(self, slide: Slide) -> str:
        """
        Extract additional information using vision models (VLM).
        Useful for image-heavy slides with diagrams, charts, etc.
        
        Args:
            slide: Slide with image_data
            
        Returns:
            Enhanced description of slide content
        """
        if not slide.image_data:
            return slide.content
        
        # TODO: Implement vision model integration
        # This would call GPT-4 Vision, Claude 3, or Gemini Vision
        # For now, return original content
        return slide.content
